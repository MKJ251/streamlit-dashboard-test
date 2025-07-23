import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

def show_download_tab(filtered_df):
    st.header("📤 Download Data")
    preview_df = filtered_df.head(20).copy()
    num_cols = preview_df.select_dtypes(include=['number']).columns
    styler = preview_df.style.set_table_styles([
        {'selector': 'thead', 'props': [('background-color', '#003366'), ('color', 'white')]},
        {'selector': 'tbody tr:hover', 'props': [('background-color', '#bfd7ff')]}
    ])
    if not num_cols.empty:
        styler = styler.highlight_max(subset=num_cols, axis=0)
    st.dataframe(styler, use_container_width=True)

    # --- Download CSV ---
    csv = filtered_df.to_csv(index=False).encode('utf-8')
    st.download_button(
        "Download CSV",
        data=csv,
        file_name="filtered_data.csv",
        mime="text/csv",
        key="download-csv-main"
    )

    # --- PowerPoint Export ---
    def generate_ppt(df):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Mindmetric Logistics Dashboard"
        subtitle.text = "Exported Report – Powered by Streamlit"

        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Key Performance Indicators"
        kpi_text = (
            f"Total Revenue: ${df['revenue_total'].sum()/1_000_000:.2f} Mn\n"
            f"Total Profit: ${df['profit'].sum()/1_000_000:.2f} Mn\n"
            f"Repeat Rate: {df['repeat_purchase_flag'].mean() * 100:.1f}%\n"
            f"ROAS Avg: {df['roas'].mean():.2f}"
        )
        textbox = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = textbox.text_frame
        tf.text = kpi_text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Sample Data Preview"
        table_df = df.head(10).reset_index(drop=True)
        rows, cols = table_df.shape
        table_shape = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        table = table_shape.table
        for i, col_name in enumerate(table_df.columns):
            table.cell(0, i).text = str(col_name)
        for row in range(rows):
            for col in range(cols):
                table.cell(row + 1, col).text = str(table_df.iloc[row, col])

        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes

    ppt_data = generate_ppt(filtered_df)
    st.download_button(
        "📥 Download PowerPoint",
        data=ppt_data,
        file_name="logistics_dashboard_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="download-ppt"
    )
