import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image

# --- Page Configuration ---
st.set_page_config(page_title="Logistics MMM Dashboard", layout="wide")

# --- Load Data ---
df = pd.read_csv("logistics_mmm_extended_data.csv")
df['week'] = pd.to_datetime(df['week'])

# --- Derived Metrics ---
df["profit"] = (df["revenue_total"] * df["profit_margin"] / 100).round(2)
df["revenue_per_order"] = (df["revenue_total"] / df["order_count"]).round(2)
df["profit_per_order"] = (df["profit"] / df["order_count"]).round(2)
df["cpc"] = (df["campaign_cost"] / df["leads_generated"]).round(2)
df["roas"] = ((df["conversions"] * df["avg_transaction_value"]) / df["campaign_cost"]).round(2)
df["conversion_rate"] = (df["conversions"] / df["leads_generated"]).round(3)

# --- Sidebar Branding and Filters ---
logo = Image.open("mindmetric_logo.png")
logo = logo.resize((100, int(logo.height * 100 / logo.width)))
st.sidebar.image(logo)

st.sidebar.markdown(
    """
    <div style="font-family: 'Segoe UI'; color:#003366; font-weight:600; font-size:16px; margin-top:10px;">
    Turning Logistics Insights into ROI 🚀
    </div>
    """,
    unsafe_allow_html=True
)
st.sidebar.markdown("---")

with st.sidebar.expander("Filter Options", expanded=True):
    date_range = st.date_input(
        "Select Date Range",
        [df['week'].min(), df['week'].max()],
        help="Filter data by date range"
    )
    regions = st.multiselect(
        "Select Region",
        df['region'].unique(),
        default=list(df['region'].unique()),
        help="Select regions to include"
    )
    customer_types = st.multiselect(
        "Select Customer Type",
        df['customer_type'].unique(),
        default=list(df['customer_type'].unique()),
        help="Select customer types"
    )

# --- Filter Data ---
start_date = pd.to_datetime(date_range[0])
end_date = pd.to_datetime(date_range[1])
filtered_df = df[
    (df['week'] >= start_date) &
    (df['week'] <= end_date) &
    (df['region'].isin(regions)) &
    (df['customer_type'].isin(customer_types))
]

# --- Top Navigation Tabs ---
tabs = st.tabs([
    "📈 Revenue & Profitability",
    "🎯 Campaign Performance",
    "🚚 Delivery & Service",
    "📣 Brand & Incidents",
    "📤 Download Data"
])

# --- KPI Cards with polished styling ---
def show_kpi_cards():
    total_revenue = filtered_df['revenue_total'].sum() / 1_000_000
    total_profit = filtered_df['profit'].sum() / 1_000_000
    repeat_rate = filtered_df['repeat_purchase_flag'].mean() * 100
    roas_avg = filtered_df['roas'].mean()

    st.markdown("""
    <style>
    .kpi-card {
        background: linear-gradient(135deg, #d0e6f7, #a3c4f3);
        border-radius: 12px;
        padding: 20px 15px;
        box-shadow: 0 4px 15px rgba(19, 43, 70, 0.2);
        text-align: center;
        font-weight: 700;
        color: #003366;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        transition: transform 0.2s ease-in-out;
        user-select: none;
    }
    .kpi-card:hover {
        transform: translateY(-6px);
        box-shadow: 0 8px 20px rgba(19, 43, 70, 0.3);
    }
    .kpi-icon {
        font-size: 34px;
        margin-bottom: 8px;
    }
    .kpi-value {
        font-size: 30px;
        margin-top: 5px;
        color: #001f4d;
    }
    .kpi-label {
        font-size: 16px;
        color: #004080;
    }
    </style>
    """, unsafe_allow_html=True)

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>💰</div>
                <div class='kpi-value'>{total_revenue:.2f} Mn</div>
                <div class='kpi-label'>Total Revenue</div>
                </div>""",
            unsafe_allow_html=True
        )
    with col2:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>📈</div>
                <div class='kpi-value'>{total_profit:.2f} Mn</div>
                <div class='kpi-label'>Total Profit</div>
                </div>""",
            unsafe_allow_html=True
        )
    with col3:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>♻️</div>
                <div class='kpi-value'>{repeat_rate:.1f}%</div>
                <div class='kpi-label'>Repeat Rate</div>
                </div>""",
            unsafe_allow_html=True
        )
    with col4:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>🎯</div>
                <div class='kpi-value'>{roas_avg:.2f}</div>
                <div class='kpi-label'>ROAS Avg</div>
                </div>""",
            unsafe_allow_html=True
        )

# Color palettes - use more saturated, professional qualitative colors
QUALITATIVE_DARK = px.colors.qualitative.Dark24
QUALITATIVE_BOLD = px.colors.qualitative.Bold
SEQ_VIRIDIS = px.colors.sequential.Viridis

# --- Tab 1: Revenue & Profitability ---
with tabs[0]:
    # st.header("📈 Revenue & Profitability Analysis")
    show_kpi_cards()
    st.write("")

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Revenue & Profit by Region and Customer Type")
        rev_df = filtered_df.groupby(["region", "customer_type"])[["revenue_total", "profit"]].mean().reset_index()
        fig1 = px.bar(
            rev_df, x="region", y="revenue_total", color="customer_type", barmode="group",
            color_discrete_sequence=QUALITATIVE_DARK,
            title=""
        )
        fig1.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig1, use_container_width=True)
    with col2:
        st.subheader("Repeat Purchase Behavior by Region")
        repeat_df = filtered_df.groupby("region")["repeat_purchase_flag"].mean().reset_index()
        fig2 = px.bar(
            repeat_df, x="region", y="repeat_purchase_flag",
            color_discrete_sequence=QUALITATIVE_DARK,
            title=""
        )
        fig2.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig2, use_container_width=True)

    st.subheader("Revenue Trends by Region")
    fig3 = px.line(
        filtered_df, x="week", y="revenue_total", color="region",
        color_discrete_sequence=QUALITATIVE_BOLD,
        title=""
    )
    fig3.update_layout(template="plotly_white", title_text="")
    st.plotly_chart(fig3, use_container_width=True)

# --- Tab 2: Campaign Performance ---
with tabs[1]:
    # st.header("🎯 Campaign Performance Overview")
    show_kpi_cards()
    st.write("")

    tab1, tab2 = st.tabs(["📊 Channel Summary", "📈 ROAS vs CAC"])
    with tab1:
        st.subheader("Lead-to-Conversion by Channel")
        conv_df = filtered_df.groupby("campaign_channel")[
            ["leads_generated", "conversions", "campaign_cost", "cpc", "roas", "customer_acquisition_cost"]
        ].mean().reset_index()
        st.dataframe(conv_df.round(2), use_container_width=True)

    with tab2:
        fig4 = px.scatter(
            conv_df, x="customer_acquisition_cost", y="roas", color="campaign_channel",
            size="conversions", hover_name="campaign_channel",
            color_discrete_sequence=QUALITATIVE_BOLD,
            title=""
        )
        fig4.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig4, use_container_width=True)

    st.subheader("Campaign Spend vs App Downloads")
    fig5 = px.scatter(
        filtered_df, x="campaign_cost", y="app_downloads", color="campaign_channel",
        size="conversions",
        color_discrete_sequence=QUALITATIVE_DARK,
        title=""
    )
    fig5.update_layout(template="plotly_white", title_text="")
    st.plotly_chart(fig5, use_container_width=True)

# --- Tab 3: Delivery & Service ---
with tabs[2]:
    # st.header("🚚 Delivery & Service Performance")
    show_kpi_cards()
    st.write("")

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Delivery Status Distribution")
        status_counts = filtered_df["delivery_status"].value_counts().reset_index(name='count')
        status_counts.columns = ['delivery_status', 'count']
        fig6 = px.pie(
            status_counts, names="delivery_status", values="count",
            color_discrete_sequence=QUALITATIVE_BOLD,
            title=""
        )
        fig6.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig6, use_container_width=True)
    with col2:
        st.subheader("Delay Reasons")
        delay_df = filtered_df[filtered_df["delivery_status"] == "Delayed"]
        delay_counts = delay_df["delay_reason"].value_counts().reset_index(name='count')
        delay_counts.columns = ['delay_reason', 'count']
        fig7 = px.bar(
            delay_counts, x="delay_reason", y="count",
            color_discrete_sequence=QUALITATIVE_DARK,
            title=""
        )
        fig7.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig7, use_container_width=True)

    st.subheader("Customer Satisfaction Trend")
    fig8 = px.line(
        filtered_df, x="week", y="customer_satisfaction_score", color="region",
        color_discrete_sequence=QUALITATIVE_BOLD,
        title=""
    )
    fig8.update_layout(template="plotly_white", title_text="")
    st.plotly_chart(fig8, use_container_width=True)

# --- Tab 4: Brand & Incidents ---
with tabs[3]:
    # st.header("📣 Brand Visibility & Incidents")
    show_kpi_cards()
    st.write("")
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Mentions & Sentiment by Channel")
        brand_df = filtered_df.groupby("media_channel")[
            ["mentions_count", "sentiment_score", "engagement_rate"]
        ].mean().reset_index()
        fig9 = px.bar(
            brand_df, x="media_channel", y="mentions_count", color="sentiment_score",
            color_continuous_scale=SEQ_VIRIDIS,
            title=""
        )
        fig9.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig9, use_container_width=True)
    with col2:
        st.subheader("Shipment Affected by Incident Type")
        inc_df = filtered_df.groupby("incident_type")["shipment_affected_count"].sum().reset_index()
        fig10 = px.bar(
            inc_df, x="incident_type", y="shipment_affected_count",
            color_discrete_sequence=QUALITATIVE_DARK,
            title=""
        )
        fig10.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig10, use_container_width=True)

# --- Tab 5: Download ---
with tabs[4]:
    # st.header("📤 Export Filtered Data")
    preview_df = filtered_df.head(20).copy()
    num_cols = preview_df.select_dtypes(include=['number']).columns

    styler = preview_df.style.set_table_styles([
        {'selector': 'thead', 'props': [('background-color', '#003366'), ('color', 'white')]},
        {'selector': 'tbody tr:hover', 'props': [('background-color', '#bfd7ff')]}
    ])

    if not num_cols.empty:
        styler = styler.highlight_max(subset=num_cols, axis=0)

    st.dataframe(styler, use_container_width=True)

    csv = filtered_df.to_csv(index=False).encode('utf-8')
    st.download_button("Download CSV", data=csv, file_name="filtered_data.csv", mime="text/csv")

    # Removed PDF export section as per request

    # PowerPoint Export (kept)
    def generate_ppt(df):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Mindmetric Logistics Dashboard"
        subtitle.text = "Exported Report – Powered by Streamlit"

        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Key Performance Indicators"
        kpi_text = (
            f"Total Revenue: ${df['revenue_total'].sum()/1_000_000:.2f} Mn\n"
            f"Total Profit: ${df['profit'].sum()/1_000_000:.2f} Mn\n"
            f"Repeat Rate: {df['repeat_purchase_flag'].mean() * 100:.1f}%\n"
            f"ROAS Avg: {df['roas'].mean():.2f}"
        )
        textbox = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = textbox.text_frame
        tf.text = kpi_text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Sample Data Preview"
        table_df = df.head(10).reset_index(drop=True)
        rows, cols = table_df.shape
        table_shape = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        table = table_shape.table
        for i, col_name in enumerate(table_df.columns):
            table.cell(0, i).text = col_name
        for row in range(rows):
            for col in range(cols):
                table.cell(row + 1, col).text = str(table_df.iloc[row, col])

        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes

    ppt_data = generate_ppt(filtered_df)
    st.download_button(
        "📥 Download PowerPoint",
        data=ppt_data,
        file_name="logistics_dashboard_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
