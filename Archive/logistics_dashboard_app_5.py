import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image

# --- Page Configuration ---
st.set_page_config(page_title="Logistics MMM Dashboard", layout="wide")

# --- Load Data ---
df = pd.read_csv("logistics_mmm_extended_data.csv")
df['week'] = pd.to_datetime(df['week'])

# --- Derived Metrics ---
df["profit"] = (df["revenue_total"] * df["profit_margin"] / 100).round(2)
df["revenue_per_order"] = (df["revenue_total"] / df["order_count"]).round(2)
df["profit_per_order"] = (df["profit"] / df["order_count"]).round(2)
df["cpc"] = (df["campaign_cost"] / df["leads_generated"]).round(2)
df["roas"] = ((df["conversions"] * df["avg_transaction_value"]) / df["campaign_cost"]).round(2)
df["conversion_rate"] = (df["conversions"] / df["leads_generated"]).round(3)

# --- Sidebar Branding and Executive Filters ---
logo = Image.open("mindmetric_logo.png")
logo = logo.resize((100, int(logo.height * 100 / logo.width)))
st.sidebar.image(logo)

st.sidebar.markdown(
    "<div style='font-family:Segoe UI;color:#003366;font-weight:600;font-size:16px;margin-top:10px;'>"
    "Turning Logistics Insights into ROI 🚀"
    "</div>", unsafe_allow_html=True
)
st.sidebar.markdown("---")

with st.sidebar.expander("Filter Options", expanded=True):
    date_range = st.date_input(
        "Select Date Range",
        [df['week'].min(), df['week'].max()],
        help="Filter data by date range"
    )
    regions = st.multiselect("Select Region", df['region'].unique(), default=list(df['region'].unique()))
    customer_types = st.multiselect("Select Customer Type", df['customer_type'].unique(), default=list(df['customer_type'].unique()))
    delivery_modes = st.multiselect("Select Delivery Mode", df['delivery_mode'].unique(), default=list(df['delivery_mode'].unique()))
    package_weight_classes = st.multiselect("Select Package Weight Class", df['package_weight_class'].unique(), default=list(df['package_weight_class'].unique()))
    service_channels = st.multiselect("Select Service Channel", df['service_channel'].unique(), default=list(df['service_channel'].unique()))
    account_types = st.multiselect("Select Account Type", df['account_type'].unique(), default=list(df['account_type'].unique()))
    customer_tiers = st.multiselect("Select Customer Tier", df['customer_tier'].unique(), default=list(df['customer_tier'].unique()))

# --- Filter Data ---
start_date = pd.to_datetime(date_range[0])
end_date = pd.to_datetime(date_range[1])
filtered_df = df[
    (df['week'] >= start_date) &
    (df['week'] <= end_date) &
    (df['region'].isin(regions)) &
    (df['customer_type'].isin(customer_types)) &
    (df['delivery_mode'].isin(delivery_modes)) &
    (df['package_weight_class'].isin(package_weight_classes)) &
    (df['service_channel'].isin(service_channels)) &
    (df['account_type'].isin(account_types)) &
    (df['customer_tier'].isin(customer_tiers))
]

tabs = st.tabs([
    "📈 Revenue & Profitability",
    "🎯 Campaign Performance",
    "🚚 Delivery & Service",
    "📣 Brand & Incidents",
    "📤 Download Data"
])

QUALITATIVE_DARK = px.colors.qualitative.Dark24
QUALITATIVE_BOLD = px.colors.qualitative.Bold
SEQ_VIRIDIS = px.colors.sequential.Viridis

def show_kpi_cards():
    total_revenue = filtered_df['revenue_total'].sum() / 1_000_000
    total_profit = filtered_df['profit'].sum() / 1_000_000
    repeat_rate = filtered_df['repeat_purchase_flag'].mean() * 100
    roas_avg = filtered_df['roas'].mean()

    st.markdown("""
    <style>
    .kpi-card {
        background: linear-gradient(135deg, #d0e6f7, #a3c4f3);
        border-radius: 12px;
        padding: 20px 15px;
        box-shadow: 0 4px 15px rgba(19, 43, 70, 0.2);
        text-align: center;
        font-weight: 700;
        color: #003366;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        transition: transform 0.2s ease-in-out;
        user-select: none;
    }
    .kpi-card:hover {
        transform: translateY(-6px);
        box-shadow: 0 8px 20px rgba(19, 43, 70, 0.3);
    }
    .kpi-icon {
        font-size: 34px;
        margin-bottom: 8px;
    }
    .kpi-value {
        font-size: 30px;
        margin-top: 5px;
        color: #001f4d;
    }
    .kpi-label {
        font-size: 16px;
        color: #004080;
    }
    </style>
    """, unsafe_allow_html=True)

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>💰</div>
                <div class='kpi-value'>{total_revenue:.2f} Mn</div>
                <div class='kpi-label'>Total Revenue</div>
                </div>""",
            unsafe_allow_html=True
        )
    with col2:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>📈</div>
                <div class='kpi-value'>{total_profit:.2f} Mn</div>
                <div class='kpi-label'>Total Profit</div>
                </div>""",
            unsafe_allow_html=True
        )
    with col3:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>♻️</div>
                <div class='kpi-value'>{repeat_rate:.1f}%</div>
                <div class='kpi-label'>Repeat Rate</div>
                </div>""",
            unsafe_allow_html=True
        )
    with col4:
        st.markdown(
            f"""<div class='kpi-card'>
                <div class='kpi-icon'>🎯</div>
                <div class='kpi-value'>{roas_avg:.2f}</div>
                <div class='kpi-label'>ROAS Avg</div>
                </div>""",
            unsafe_allow_html=True
        )

# --- Define KPI Cards with YoY Growth Indicators ---
def show_kpi_cards_with_yoy():
    current_year = filtered_df['week'].dt.year.max()
    previous_year = current_year - 1
    curr_year_df = filtered_df[filtered_df['week'].dt.year == current_year]
    prev_year_df = filtered_df[filtered_df['week'].dt.year == previous_year]
    curr_vals = curr_year_df.agg({
        'revenue_total': 'sum',
        'profit': 'sum',
        'repeat_purchase_flag': 'mean',
        'roas': 'mean'
    })
    prev_vals = prev_year_df.agg({
        'revenue_total': 'sum',
        'profit': 'sum',
        'repeat_purchase_flag': 'mean',
        'roas': 'mean'
    })
    kpi_yoy = ((curr_vals - prev_vals) / prev_vals) * 100

    total_revenue = curr_vals['revenue_total'] / 1_000_000
    total_profit = curr_vals['profit'] / 1_000_000
    repeat_rate = curr_vals['repeat_purchase_flag'] * 100
    roas_avg = curr_vals['roas']

    def arrow(val):
        if pd.isna(val):
            return ""
        elif val > 0:
            return f' <span style="color:#1AA83B;font-size:15px;">&#9650; {val:.1f}%</span>'
        elif val < 0:
            return f' <span style="color:#E02424;font-size:15px;">&#9660; {abs(val):.1f}%</span>'
        else:
            return ' <span style="color:#888888;font-size:15px;">0.0%</span>'

    st.markdown(
        """
        <style>
        .kpi-card{background:linear-gradient(135deg,#d0e6f7,#a3c4f3);border-radius:12px;padding:20px 15px;box-shadow:0 4px 15px rgba(19,43,70,0.2);text-align:center;font-weight:700;color:#003366;font-family:'Segoe UI',Tahoma,Geneva,Verdana,sans-serif;user-select:none;}.kpi-value{font-size:30px;margin-top:5px;color:#001f4d;}.kpi-label{font-size:16px;color:#004080;}.kpi-delta{font-size:15px;font-weight:500;margin-top:2px;}
        </style>
        """, unsafe_allow_html=True)

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown(
            f"""<div class='kpi-card'>💰<div class='kpi-value'>{total_revenue:.2f} Mn</div>
                <div class='kpi-label'>Total Revenue</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['revenue_total'])}</div></div>""",
            unsafe_allow_html=True)
    with col2:
        st.markdown(
            f"""<div class='kpi-card'>📈<div class='kpi-value'>{total_profit:.2f} Mn</div>
                <div class='kpi-label'>Total Profit</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['profit'])}</div></div>""",
            unsafe_allow_html=True)
    with col3:
        st.markdown(
            f"""<div class='kpi-card'>♻️<div class='kpi-value'>{repeat_rate:.1f}%</div>
                <div class='kpi-label'>Repeat Rate</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['repeat_purchase_flag'])}</div></div>""",
            unsafe_allow_html=True)
    with col4:
        st.markdown(
            f"""<div class='kpi-card'>🎯<div class='kpi-value'>{roas_avg:.2f}</div>
                <div class='kpi-label'>ROAS Avg</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['roas'])}</div></div>""",
            unsafe_allow_html=True)

# --- TAB 1: Revenue & Profitability (Executive Enhanced) ---
with tabs[0]:
    st.markdown("### Executive Summary")
    show_kpi_cards_with_yoy()
    st.write("")
    st.markdown("---")

    # 1. Revenue by Region and Revenue by Customer Type (side by side)
    col1, col2 = st.columns(2)
    with col1:
        st.markdown("**Revenue by Region**")
        rev_by_region = filtered_df.groupby('region')['revenue_total'].sum().reset_index()
        fig_region = px.bar(
            rev_by_region, x='region', y='revenue_total', color='region',
            color_discrete_sequence=QUALITATIVE_DARK, title=None
        )
        fig_region.update_layout(template="plotly_white", xaxis_title="Region", yaxis_title="Total Revenue")
        st.plotly_chart(fig_region, use_container_width=True)
    with col2:
        st.markdown("**Revenue by Customer Type (B2B vs B2C)**")
        rev_by_custtype = filtered_df.groupby('customer_type')['revenue_total'].sum().reset_index()
        fig_custtype = px.bar(
            rev_by_custtype, x='customer_type', y='revenue_total', color='customer_type',
            color_discrete_sequence=QUALITATIVE_BOLD, title=None
        )
        fig_custtype.update_layout(template="plotly_white", xaxis_title="Customer Type", yaxis_title="Total Revenue")
        st.plotly_chart(fig_custtype, use_container_width=True)

    # 2. Revenue Trend by Region (single chart)
    st.markdown("**Revenue Trend by Region**")
    rev_trend_region = filtered_df.groupby(['week', 'region'])['revenue_total'].sum().reset_index()
    fig_trend_region = px.line(
        rev_trend_region, x='week', y='revenue_total', color='region',
        color_discrete_sequence=QUALITATIVE_DARK
    )
    fig_trend_region.update_layout(template="plotly_white", xaxis_title="Week", yaxis_title="Revenue")
    st.plotly_chart(fig_trend_region, use_container_width=True)

    # 3. Revenue by Delivery Mode and Package Weight Class (side by side)
    col3, col4 = st.columns(2)
    with col3:
        st.markdown("**Revenue by Delivery Mode**")
        rev_by_mode = filtered_df.groupby('delivery_mode')['revenue_total'].sum().reset_index()
        fig_mode = px.bar(
            rev_by_mode, x='delivery_mode', y='revenue_total', color='delivery_mode',
            color_discrete_sequence=QUALITATIVE_DARK
        )
        fig_mode.update_layout(template="plotly_white", xaxis_title="Delivery Mode", yaxis_title="Revenue")
        st.plotly_chart(fig_mode, use_container_width=True)
    with col4:
        st.markdown("**Revenue by Package Weight Class**")
        rev_by_pkg = filtered_df.groupby('package_weight_class')['revenue_total'].sum().reset_index()
        fig_pkg = px.bar(
            rev_by_pkg, x='package_weight_class', y='revenue_total', color='package_weight_class',
            color_discrete_sequence=QUALITATIVE_BOLD
        )
        fig_pkg.update_layout(template="plotly_white", xaxis_title="Weight Class", yaxis_title="Revenue")
        st.plotly_chart(fig_pkg, use_container_width=True)

    # 4. Revenue by Service Channel, Account Type, Customer Tier as pie charts side by side
    col5, col6, col7 = st.columns(3)
    with col5:
        st.markdown("**Revenue by Service Channel**")
        pie_service = filtered_df.groupby('service_channel')['revenue_total'].sum().reset_index()
        fig_service = px.pie(
            pie_service, names='service_channel', values='revenue_total',
            color_discrete_sequence=QUALITATIVE_DARK
        )
        fig_service.update_traces(textinfo='percent+label')
        st.plotly_chart(fig_service, use_container_width=True)
    with col6:
        st.markdown("**Revenue by Account Type**")
        pie_account = filtered_df.groupby('account_type')['revenue_total'].sum().reset_index()
        fig_account = px.pie(
            pie_account, names='account_type', values='revenue_total',
            color_discrete_sequence=QUALITATIVE_BOLD
        )
        fig_account.update_traces(textinfo='percent+label')
        st.plotly_chart(fig_account, use_container_width=True)
    with col7:
        st.markdown("**Revenue by Customer Tier**")
        pie_tier = filtered_df.groupby('customer_tier')['revenue_total'].sum().reset_index()
        fig_tier = px.pie(
            pie_tier, names='customer_tier', values='revenue_total',
            color_discrete_sequence=QUALITATIVE_DARK
        )
        fig_tier.update_traces(textinfo='percent+label')
        st.plotly_chart(fig_tier, use_container_width=True)

    # 5. Trend charts at the bottom: Weekly CAC and Churn Rate (single charts)
    st.markdown("---")
    st.subheader("Customer Metrics Trends (Weekly)")
    st.markdown("**Weekly Customer Acquisition Cost**")
    cac_trend = filtered_df.groupby('week')['customer_acquisition_cost'].mean().reset_index()
    fig_cac = px.line(
        cac_trend, x='week', y='customer_acquisition_cost',
        color_discrete_sequence=QUALITATIVE_DARK, labels={"customer_acquisition_cost": "Avg Acquisition Cost"}
    )
    fig_cac.update_layout(template="plotly_white", xaxis_title="Week", yaxis_title="Acquisition Cost")
    st.plotly_chart(fig_cac, use_container_width=True)

    st.markdown("**Weekly Customer Churn Rate**")
    churn_trend = filtered_df.groupby('week')['customer_churn_rate'].mean().reset_index()
    fig_churn = px.line(
        churn_trend, x='week', y='customer_churn_rate',
        color_discrete_sequence=QUALITATIVE_BOLD, labels={"customer_churn_rate": "Avg Churn Rate"}
    )
    fig_churn.update_layout(template="plotly_white", xaxis_title="Week", yaxis_title="Churn Rate")
    st.plotly_chart(fig_churn, use_container_width=True)

# --- Other Tabs (not shown for brevity, but keep as before) ---

# --- Tab 2: Campaign Performance ---
with tabs[1]:
    # st.header("🎯 Campaign Performance Overview")
    show_kpi_cards()
    st.write("")

    tab1, tab2 = st.tabs(["📊 Channel Summary", "📈 ROAS vs CAC"])
    with tab1:
        st.subheader("Lead-to-Conversion by Channel")
        conv_df = filtered_df.groupby("campaign_channel")[
            ["leads_generated", "conversions", "campaign_cost", "cpc", "roas", "customer_acquisition_cost"]
        ].mean().reset_index()
        st.dataframe(conv_df.round(2), use_container_width=True)

    with tab2:
        fig4 = px.scatter(
            conv_df, x="customer_acquisition_cost", y="roas", color="campaign_channel",
            size="conversions", hover_name="campaign_channel",
            color_discrete_sequence=QUALITATIVE_BOLD,
            title=""
        )
        fig4.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig4, use_container_width=True)

    st.subheader("Campaign Spend vs App Downloads")
    fig5 = px.scatter(
        filtered_df, x="campaign_cost", y="app_downloads", color="campaign_channel",
        size="conversions",
        color_discrete_sequence=QUALITATIVE_DARK,
        title=""
    )
    fig5.update_layout(template="plotly_white", title_text="")
    st.plotly_chart(fig5, use_container_width=True)

# --- Tab 3: Delivery & Service ---
with tabs[2]:
    # st.header("🚚 Delivery & Service Performance")
    show_kpi_cards()
    st.write("")

    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Delivery Status Distribution")
        status_counts = filtered_df["delivery_status"].value_counts().reset_index(name='count')
        status_counts.columns = ['delivery_status', 'count']
        fig6 = px.pie(
            status_counts, names="delivery_status", values="count",
            color_discrete_sequence=QUALITATIVE_BOLD,
            title=""
        )
        fig6.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig6, use_container_width=True)
    with col2:
        st.subheader("Delay Reasons")
        delay_df = filtered_df[filtered_df["delivery_status"] == "Delayed"]
        delay_counts = delay_df["delay_reason"].value_counts().reset_index(name='count')
        delay_counts.columns = ['delay_reason', 'count']
        fig7 = px.bar(
            delay_counts, x="delay_reason", y="count",
            color_discrete_sequence=QUALITATIVE_DARK,
            title=""
        )
        fig7.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig7, use_container_width=True)

    st.subheader("Customer Satisfaction Trend")
    fig8 = px.line(
        filtered_df, x="week", y="customer_satisfaction_score", color="region",
        color_discrete_sequence=QUALITATIVE_BOLD,
        title=""
    )
    fig8.update_layout(template="plotly_white", title_text="")
    st.plotly_chart(fig8, use_container_width=True)

# --- Tab 4: Brand & Incidents ---
with tabs[3]:
    # st.header("📣 Brand Visibility & Incidents")
    show_kpi_cards()
    st.write("")
    col1, col2 = st.columns(2)
    with col1:
        st.subheader("Mentions & Sentiment by Channel")
        brand_df = filtered_df.groupby("media_channel")[
            ["mentions_count", "sentiment_score", "engagement_rate"]
        ].mean().reset_index()
        fig9 = px.bar(
            brand_df, x="media_channel", y="mentions_count", color="sentiment_score",
            color_continuous_scale=SEQ_VIRIDIS,
            title=""
        )
        fig9.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig9, use_container_width=True)
    with col2:
        st.subheader("Shipment Affected by Incident Type")
        inc_df = filtered_df.groupby("incident_type")["shipment_affected_count"].sum().reset_index()
        fig10 = px.bar(
            inc_df, x="incident_type", y="shipment_affected_count",
            color_discrete_sequence=QUALITATIVE_DARK,
            title=""
        )
        fig10.update_layout(template="plotly_white", title_text="")
        st.plotly_chart(fig10, use_container_width=True)

# --- Tab 5: Download ---
with tabs[4]:
    # st.header("📤 Export Filtered Data")
    preview_df = filtered_df.head(20).copy()
    num_cols = preview_df.select_dtypes(include=['number']).columns

    styler = preview_df.style.set_table_styles([
        {'selector': 'thead', 'props': [('background-color', '#003366'), ('color', 'white')]},
        {'selector': 'tbody tr:hover', 'props': [('background-color', '#bfd7ff')]}
    ])

    if not num_cols.empty:
        styler = styler.highlight_max(subset=num_cols, axis=0)

    st.dataframe(styler, use_container_width=True)

    csv = filtered_df.to_csv(index=False).encode('utf-8')
    st.download_button("Download CSV", data=csv, file_name="filtered_data.csv", mime="text/csv")


# --- Download Data Tab ---
with tabs[4]:
    preview_df = filtered_df.head(20).copy()
    num_cols = preview_df.select_dtypes(include=['number']).columns
    styler = preview_df.style.set_table_styles([
        {'selector': 'thead', 'props': [('background-color', '#003366'), ('color', 'white')]},
        {'selector': 'tbody tr:hover', 'props': [('background-color', '#bfd7ff')]}
    ])
    if not num_cols.empty:
        styler = styler.highlight_max(subset=num_cols, axis=0)
    st.dataframe(styler, use_container_width=True)

    csv = filtered_df.to_csv(index=False).encode('utf-8')
    st.download_button("Download CSV", data=csv, file_name="filtered_data.csv", mime="text/csv")

    # PowerPoint Export
    def generate_ppt(df):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Mindmetric Logistics Dashboard"
        subtitle.text = "Exported Report – Powered by Streamlit"
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Key Performance Indicators"
        kpi_text = (
            f"Total Revenue: ${df['revenue_total'].sum()/1_000_000:.2f} Mn\n"
            f"Total Profit: ${df['profit'].sum()/1_000_000:.2f} Mn\n"
            f"Repeat Rate: {df['repeat_purchase_flag'].mean() * 100:.1f}%\n"
            f"ROAS Avg: {df['roas'].mean():.2f}"
        )
        textbox = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = textbox.text_frame
        tf.text = kpi_text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Sample Data Preview"
        table_df = df.head(10).reset_index(drop=True)
        rows, cols = table_df.shape
        table_shape = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        table = table_shape.table
        for i, col_name in enumerate(table_df.columns):
            table.cell(0, i).text = col_name
        for row in range(rows):
            for col in range(cols):
                table.cell(row + 1, col).text = str(table_df.iloc[row, col])

        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes

    ppt_data = generate_ppt(filtered_df)
    st.download_button(
        "📥 Download PowerPoint",
        data=ppt_data,
        file_name="logistics_dashboard_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
