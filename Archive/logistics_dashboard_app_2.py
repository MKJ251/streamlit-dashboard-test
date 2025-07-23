import streamlit as st
import pandas as pd
import plotly.express as px
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image

# --- Page Configuration ---
st.set_page_config(page_title="Logistics MMM Dashboard", layout="wide")

# --- Load Data ---
df = pd.read_csv("logistics_mmm_extended_data.csv")
df['week'] = pd.to_datetime(df['week'])

# --- Derived Metrics ---
df["profit"] = (df["revenue_total"] * df["profit_margin"] / 100).round(2)
df["revenue_per_order"] = (df["revenue_total"] / df["order_count"]).round(2)
df["profit_per_order"] = (df["profit"] / df["order_count"]).round(2)
df["cpc"] = (df["campaign_cost"] / df["leads_generated"]).round(2)
df["roas"] = ((df["conversions"] * df["avg_transaction_value"]) / df["campaign_cost"]).round(2)
df["conversion_rate"] = (df["conversions"] / df["leads_generated"]).round(3)

# --- Sidebar Branding and Executive Filters ---
logo = Image.open("mindmetric_logo.png")
logo = logo.resize((100, int(logo.height * 100 / logo.width)))
st.sidebar.image(logo)

st.sidebar.markdown(
    "<div style='font-family:Segoe UI;color:#003366;font-weight:600;font-size:16px;margin-top:10px;'>"
    "Turning Logistics Insights into ROI 🚀"
    "</div>", unsafe_allow_html=True
)
st.sidebar.markdown("---")

with st.sidebar.expander("Filter Options", expanded=True):
    date_range = st.date_input(
        "Select Date Range",
        [df['week'].min(), df['week'].max()],
        help="Filter data by date range"
    )
    regions = st.multiselect("Select Region", df['region'].unique(), default=list(df['region'].unique()))
    customer_types = st.multiselect("Select Customer Type", df['customer_type'].unique(), default=list(df['customer_type'].unique()))
    delivery_modes = st.multiselect("Select Delivery Mode", df['delivery_mode'].unique(), default=list(df['delivery_mode'].unique()))
    package_weight_classes = st.multiselect("Select Package Weight Class", df['package_weight_class'].unique(), default=list(df['package_weight_class'].unique()))
    service_channels = st.multiselect("Select Service Channel", df['service_channel'].unique(), default=list(df['service_channel'].unique()))
    account_types = st.multiselect("Select Account Type", df['account_type'].unique(), default=list(df['account_type'].unique()))
    customer_tiers = st.multiselect("Select Customer Tier", df['customer_tier'].unique(), default=list(df['customer_tier'].unique()))

# --- Filter Data ---
start_date = pd.to_datetime(date_range[0])
end_date = pd.to_datetime(date_range[1])
filtered_df = df[
    (df['week'] >= start_date) &
    (df['week'] <= end_date) &
    (df['region'].isin(regions)) &
    (df['customer_type'].isin(customer_types)) &
    (df['delivery_mode'].isin(delivery_modes)) &
    (df['package_weight_class'].isin(package_weight_classes)) &
    (df['service_channel'].isin(service_channels)) &
    (df['account_type'].isin(account_types)) &
    (df['customer_tier'].isin(customer_tiers))
]

# --- Top Navigation Tabs ---
tabs = st.tabs([
    "📈 Revenue & Profitability",
    "🎯 Campaign Performance",
    "🚚 Delivery & Service",
    "📣 Brand & Incidents",
    "📤 Download Data"
])

QUALITATIVE_DARK = px.colors.qualitative.Dark24
QUALITATIVE_BOLD = px.colors.qualitative.Bold
SEQ_VIRIDIS = px.colors.sequential.Viridis

# === Executive KPI Snapshot with YoY Deltas ===
def show_kpi_cards_with_yoy():
    current_year = filtered_df['week'].dt.year.max()
    previous_year = current_year - 1
    curr_year_df = filtered_df[filtered_df['week'].dt.year == current_year]
    prev_year_df = filtered_df[filtered_df['week'].dt.year == previous_year]

    curr_vals = curr_year_df.agg({
        'revenue_total': 'sum',
        'profit': 'sum',
        'repeat_purchase_flag': 'mean',
        'roas': 'mean'
    })
    prev_vals = prev_year_df.agg({
        'revenue_total': 'sum',
        'profit': 'sum',
        'repeat_purchase_flag': 'mean',
        'roas': 'mean'
    })
    kpi_yoy = ((curr_vals - prev_vals) / prev_vals) * 100

    total_revenue = curr_vals['revenue_total'] / 1_000_000
    total_profit = curr_vals['profit'] / 1_000_000
    repeat_rate = curr_vals['repeat_purchase_flag'] * 100
    roas_avg = curr_vals['roas']

    def arrow(val):
        if pd.isna(val):
            return ""
        elif val > 0:
            return f' <span style="color:#1AA83B;font-size:15px;">&#9650; {val:.1f}%</span>'
        elif val < 0:
            return f' <span style="color:#E02424;font-size:15px;">&#9660; {abs(val):.1f}%</span>'
        else:
            return ' <span style="color:#888888;font-size:15px;">0.0%</span>'

    st.markdown(
        """
        <style>
        .kpi-card{background:linear-gradient(135deg,#d0e6f7,#a3c4f3);border-radius:12px;padding:20px 15px;box-shadow:0 4px 15px rgba(19,43,70,0.2);text-align:center;font-weight:700;color:#003366;font-family:'Segoe UI',Tahoma,Geneva,Verdana,sans-serif;user-select:none;}.kpi-value{font-size:30px;margin-top:5px;color:#001f4d;}.kpi-label{font-size:16px;color:#004080;}.kpi-delta{font-size:15px;font-weight:500;margin-top:2px;}
        </style>
        """, unsafe_allow_html=True)

    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown(
            f"""<div class='kpi-card'>💰<div class='kpi-value'>{total_revenue:.2f} Mn</div>
                <div class='kpi-label'>Total Revenue</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['revenue_total'])}</div></div>""",
            unsafe_allow_html=True)
    with col2:
        st.markdown(
            f"""<div class='kpi-card'>📈<div class='kpi-value'>{total_profit:.2f} Mn</div>
                <div class='kpi-label'>Total Profit</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['profit'])}</div></div>""",
            unsafe_allow_html=True)
    with col3:
        st.markdown(
            f"""<div class='kpi-card'>♻️<div class='kpi-value'>{repeat_rate:.1f}%</div>
                <div class='kpi-label'>Repeat Rate</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['repeat_purchase_flag'])}</div></div>""",
            unsafe_allow_html=True)
    with col4:
        st.markdown(
            f"""<div class='kpi-card'>🎯<div class='kpi-value'>{roas_avg:.2f}</div>
                <div class='kpi-label'>ROAS Avg</div>
                <div class='kpi-delta'>{arrow(kpi_yoy['roas'])}</div></div>""",
            unsafe_allow_html=True)

# --- TAB 1: Revenue & Profitability (Executive Enhanced) ---
with tabs[0]:
    st.markdown("<h4>Executive Summary</h4>", unsafe_allow_html=True)
    show_kpi_cards_with_yoy()
    st.write("")
    st.markdown("---")
    st.subheader("Diagnostic Sales KPIs & Trends")

    # Revenue & Profit by Region and Customer Type
    rev_df = filtered_df.groupby(["region", "customer_type"])[["revenue_total", "profit"]].sum().reset_index()
    fig1 = px.bar(rev_df, x="region", y="revenue_total", color="customer_type", barmode="group", color_discrete_sequence=QUALITATIVE_DARK)
    fig1.update_layout(template="plotly_white", yaxis_title="Revenue", xaxis_title="Region")
    st.plotly_chart(fig1, use_container_width=True)

    # Revenue & Profit by Delivery Mode and Customer Type
    rev_mode_df = filtered_df.groupby(['delivery_mode', 'customer_type'])[['revenue_total', 'profit']].sum().reset_index()
    fig2 = px.bar(rev_mode_df, x="delivery_mode", y="revenue_total", color="customer_type", barmode="group", color_discrete_sequence=QUALITATIVE_BOLD)
    fig2.update_layout(template="plotly_white", yaxis_title="Revenue", xaxis_title="Delivery Mode")
    st.plotly_chart(fig2, use_container_width=True)

    # Revenue & Profit by Package Weight Class and Customer Type
    rev_weight_df = filtered_df.groupby(['package_weight_class', 'customer_type'])[['revenue_total', 'profit']].sum().reset_index()
    fig3 = px.bar(rev_weight_df, x="package_weight_class", y="revenue_total", color="customer_type", barmode="group", color_discrete_sequence=QUALITATIVE_DARK)
    fig3.update_layout(template="plotly_white", yaxis_title="Revenue", xaxis_title="Weight Class")
    st.plotly_chart(fig3, use_container_width=True)

    # ROAS & Profit Margin Trend Over Time
    trends_df = filtered_df.copy()
    trends_df["profit_margin"] = trends_df["profit_margin"].astype(float)
    trend_line = trends_df.groupby("week")[["roas", "profit_margin"]].mean().reset_index()
    fig4 = px.line(trend_line, x="week", y=["roas", "profit_margin"], color_discrete_sequence=QUALITATIVE_BOLD, labels={"value":"Value", "variable":"Metric"})
    fig4.update_layout(template="plotly_white", title_text="")
    st.plotly_chart(fig4, use_container_width=True)

#    # Top 10 High-Value Customers
#    st.subheader("Top 10 High-Value Customers")
#    top_customers = filtered_df.groupby('customer_name')[['revenue_total', 'profit', 'order_count']].sum().reset_index()
#    top_customers = top_customers.sort_values(by='revenue_total', ascending=False).head(10)
#    fig5 = px.bar(top_customers, x='customer_name', y='revenue_total', color='profit', color_continuous_scale=px.colors.sequential.Plasma, title="Top 10 Customers by Revenue")
#    fig5.update_layout(template="plotly_white", yaxis_title="Revenue", xaxis_title="Customer", showlegend=False)
#    st.plotly_chart(fig5, use_container_width=True)
#    st.dataframe(top_customers)
# 
    # Repeat vs New Revenue Split
    repeat_revenue = filtered_df.groupby('repeat_purchase_flag')['revenue_total'].sum().reset_index()
    repeat_revenue['customer_type_desc'] = repeat_revenue['repeat_purchase_flag'].map({0: 'New', 1: 'Repeat'})
    fig6 = px.pie(repeat_revenue, names='customer_type_desc', values='revenue_total', title="Revenue from Repeat vs New Customers", color_discrete_sequence=QUALITATIVE_BOLD)
    st.plotly_chart(fig6, use_container_width=True)

# --- Other Tabs ---
# Continue using filtered_df with new filters for other tabs (Campaign, Service, Brand, Download) as before.

# --- Download Data Tab ---
with tabs[4]:
    preview_df = filtered_df.head(20).copy()
    num_cols = preview_df.select_dtypes(include=['number']).columns

    styler = preview_df.style.set_table_styles([
        {'selector': 'thead', 'props': [('background-color', '#003366'), ('color', 'white')]},
        {'selector': 'tbody tr:hover', 'props': [('background-color', '#bfd7ff')]}
    ])

    if not num_cols.empty:
        styler = styler.highlight_max(subset=num_cols, axis=0)

    st.dataframe(styler, use_container_width=True)

    csv = filtered_df.to_csv(index=False).encode('utf-8')
    st.download_button("Download CSV", data=csv, file_name="filtered_data.csv", mime="text/csv")

    # PowerPoint Export
    def generate_ppt(df):
        prs = Presentation()
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Mindmetric Logistics Dashboard"
        subtitle.text = "Exported Report – Powered by Streamlit"

        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        shapes.title.text = "Key Performance Indicators"
        kpi_text = (
            f"Total Revenue: ${df['revenue_total'].sum()/1_000_000:.2f} Mn\n"
            f"Total Profit: ${df['profit'].sum()/1_000_000:.2f} Mn\n"
            f"Repeat Rate: {df['repeat_purchase_flag'].mean() * 100:.1f}%\n"
            f"ROAS Avg: {df['roas'].mean():.2f}"
        )
        textbox = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        tf = textbox.text_frame
        tf.text = kpi_text
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)

        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Sample Data Preview"
        table_df = df.head(10).reset_index(drop=True)
        rows, cols = table_df.shape
        table_shape = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(4))
        table = table_shape.table
        for i, col_name in enumerate(table_df.columns):
            table.cell(0, i).text = col_name
        for row in range(rows):
            for col in range(cols):
                table.cell(row + 1, col).text = str(table_df.iloc[row, col])

        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes

    ppt_data = generate_ppt(filtered_df)
    st.download_button(
        "📥 Download PowerPoint",
        data=ppt_data,
        file_name="logistics_dashboard_report.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )